import argparse
from pathlib import Path
from pptx import Presentation

EMU_PER_INCH = 914400
SLIDE_W = 13.333333 * EMU_PER_INCH
SLIDE_H = 7.5 * EMU_PER_INCH


def shape_has_text(shape):
    return getattr(shape, "has_text_frame", False) and bool(shape.text_frame.text.strip())


def validate(path: Path, expected_slides: int | None = None):
    issues = []
    warnings = []
    if not path.exists():
        issues.append(f"missing file: {path}")
        return issues, warnings, None
    if path.stat().st_size <= 0:
        issues.append("file is empty")
        return issues, warnings, None

    prs = Presentation(path)
    slide_count = len(prs.slides)
    if expected_slides is not None and slide_count != expected_slides:
        issues.append(f"slide count mismatch: got {slide_count}, expected {expected_slides}")

    for i, slide in enumerate(prs.slides, 1):
        shapes = list(slide.shapes)
        if not shapes:
            issues.append(f"slide {i}: empty slide")
            continue
        text_shapes = [s for s in shapes if shape_has_text(s)]
        if not text_shapes:
            warnings.append(f"slide {i}: no editable text detected")
        if len(text_shapes) > 18:
            warnings.append(f"slide {i}: many text boxes ({len(text_shapes)}), may be crowded")
        picture_count = sum(1 for s in shapes if getattr(s, 'shape_type', None) == 13)
        if picture_count == 0 and len(shapes) < 6:
            warnings.append(f"slide {i}: no picture and few visual shapes; may feel text-only")
        for s in shapes:
            if s.left < -10000 or s.top < -10000 or s.left + s.width > SLIDE_W + 10000 or s.top + s.height > SLIDE_H + 10000:
                warnings.append(f"slide {i}: shape may overflow bounds")
            if shape_has_text(s):
                text = s.text_frame.text.strip()
                if '�' in text or '?' in text or '????' in text:
                    issues.append(f"slide {i}: possible mojibake/replacement text: {text[:40]}")
                if len(text) > 150:
                    warnings.append(f"slide {i}: long text box ({len(text)} chars), possible overflow risk")
                for p in s.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.size and r.font.size.pt < 8:
                            warnings.append(f"slide {i}: tiny font < 8pt")
    return issues, warnings, slide_count


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--expected-slides", type=int)
    args = ap.parse_args()
    issues, warnings, slide_count = validate(Path(args.pptx), args.expected_slides)
    status = "FAIL" if issues else ("WARN" if warnings else "PASS")
    print("PPT_VALIDATE_RESULT")
    print(f"status: {status}")
    print(f"slides: {slide_count}")
    print(f"issues: {len(issues)}")
    for item in issues:
        print(f"- ISSUE: {item}")
    print(f"warnings: {len(warnings)}")
    for item in warnings:
        print(f"- WARN: {item}")
    raise SystemExit(1 if issues else 0)


if __name__ == "__main__":
    main()
