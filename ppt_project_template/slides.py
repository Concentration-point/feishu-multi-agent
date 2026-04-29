from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from layouts import set_bg, add_textbox, add_title, add_footer, add_card, add_bullet_list, BLANK
from theme import SLIDE_W, SLIDE_H, MARGIN_X, MARGIN_Y, TITLE_SIZE, SUBTITLE_SIZE, BODY_SIZE, SMALL_SIZE, rgb


def title_slide(prs, theme, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(slide, theme.bg)
    add_card(slide, Inches(0.7), Inches(0.7), Inches(12.0), Inches(6.1), theme, fill=theme.panel)
    add_textbox(slide, Inches(1.15), Inches(2.25), Inches(10.6), Inches(0.8), data["title"], theme, Pt(38), True, theme.primary, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.4), Inches(3.15), Inches(10.0), Inches(0.4), data.get("subtitle", ""), theme, SUBTITLE_SIZE, False, theme.muted, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.4), Inches(5.7), Inches(10.0), Inches(0.25), data.get("meta", ""), theme, SMALL_SIZE, False, theme.muted, PP_ALIGN.CENTER)
    return slide


def agenda_slide(prs, theme, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(slide, theme.bg)
    add_title(slide, data["title"], theme)
    items = data.get("items", [])[:5]
    y = Inches(1.65)
    for n, item in enumerate(items, 1):
        add_card(slide, Inches(1.0), y, Inches(11.2), Inches(0.72), theme)
        add_textbox(slide, Inches(1.25), y + Inches(0.16), Inches(0.55), Inches(0.28), f"{n:02d}", theme, Pt(14), True, theme.accent)
        add_textbox(slide, Inches(2.0), y + Inches(0.14), Inches(9.7), Inches(0.32), item, theme, Pt(17), True)
        y += Inches(0.9)
    add_footer(slide, theme, idx)
    return slide


def section_slide(prs, theme, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(slide, theme.primary)
    add_textbox(slide, Inches(0.9), Inches(1.0), Inches(3.0), Inches(0.8), data.get("section_no", ""), theme, Pt(42), True, theme.accent)
    add_textbox(slide, Inches(0.95), Inches(2.65), Inches(10.5), Inches(0.8), data["title"], theme, Pt(36), True, theme.bg)
    add_textbox(slide, Inches(1.0), Inches(3.55), Inches(9.5), Inches(0.4), data.get("subtitle", ""), theme, Pt(18), False, theme.bg)
    return slide


def three_cards_slide(prs, theme, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(slide, theme.bg)
    add_title(slide, data["title"], theme)
    cards = data.get("cards", [])[:3]
    card_w = Inches(3.65)
    gap = Inches(0.35)
    x0 = Inches(0.9)
    y = Inches(2.0)
    for i, card in enumerate(cards):
        x = x0 + i * (card_w + gap)
        add_card(slide, x, y, card_w, Inches(3.3), theme)
        add_textbox(slide, x + Inches(0.25), y + Inches(0.35), card_w - Inches(0.5), Inches(0.4), card.get("title", ""), theme, Pt(19), True, theme.primary)
        add_textbox(slide, x + Inches(0.25), y + Inches(1.0), card_w - Inches(0.5), Inches(1.6), card.get("body", ""), theme, Pt(13), False, theme.text)
    add_footer(slide, theme, idx)
    return slide


def process_slide(prs, theme, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(slide, theme.bg)
    add_title(slide, data["title"], theme)
    steps = data.get("steps", [])[:5]
    x = Inches(0.9)
    y = Inches(3.0)
    w = Inches(2.1)
    for i, step in enumerate(steps):
        add_card(slide, x, y, w, Inches(0.8), theme, fill=theme.primary if i == 0 else theme.panel)
        add_textbox(slide, x + Inches(0.12), y + Inches(0.24), w - Inches(0.24), Inches(0.25), step, theme, Pt(13), True, theme.bg if i == 0 else theme.text, PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_textbox(slide, x + w + Inches(0.05), y + Inches(0.24), Inches(0.35), Inches(0.25), "→", theme, Pt(18), True, theme.accent, PP_ALIGN.CENTER)
        x += Inches(2.42)
    add_footer(slide, theme, idx)
    return slide


def summary_slide(prs, theme, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(slide, theme.bg)
    add_title(slide, data["title"], theme)
    add_bullet_list(slide, Inches(1.15), Inches(1.8), Inches(10.7), Inches(3.8), data.get("points", [])[:5], theme, Pt(17))
    add_footer(slide, theme, idx)
    return slide


def qa_slide(prs, theme, data, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(slide, theme.primary)
    add_textbox(slide, Inches(1.0), Inches(2.6), Inches(11.3), Inches(0.9), data.get("title", "Q&A"), theme, Pt(54), True, theme.bg, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.0), Inches(3.65), Inches(11.3), Inches(0.35), "Questions and discussion", theme, Pt(18), False, theme.bg, PP_ALIGN.CENTER)
    return slide

BUILDERS = {
    "title": title_slide,
    "agenda": agenda_slide,
    "section": section_slide,
    "three_cards": three_cards_slide,
    "process": process_slide,
    "summary": summary_slide,
    "qa": qa_slide,
}
