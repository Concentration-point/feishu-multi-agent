from pathlib import Path
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUT = Path(__file__).resolve().parent / 'output' / 'surgical_robot_claims_presentation.pptx'
W, H = Inches(13.333333), Inches(7.5)

COL = {
    'bg': 'F7F8FA', 'panel': 'FFFFFF', 'primary': '1F4E79', 'accent': 'D89A2B',
    'text': '222222', 'muted': '6B7280', 'dark': '0F172A', 'soft': 'E8EEF5'
}
FONT = 'Microsoft YaHei'


def rgb(h):
    h = h.strip('#')
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:], 16))


def bg(slide, color='bg'):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = rgb(COL[color])


def txt(slide, x, y, w, h, s, size=16, bold=False, color='text', align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = s
    if align is not None:
        p.alignment = align
    r = p.runs[0]
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(COL[color])
    return box


def bullet(slide, x, y, w, h, items, size=15, color='text'):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(item)
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(COL[color])
        p.space_after = Pt(6)
    return box


def rect(slide, x, y, w, h, fill='panel', line=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COL[fill])
    shape.line.color.rgb = rgb(COL[line or fill])
    return shape


def title(slide, s, sub=None, idx=None):
    txt(slide, 0.65, 0.45, 10.8, 0.45, s, 24, True, 'primary')
    if sub:
        txt(slide, 0.67, 0.92, 10.8, 0.28, sub, 10, False, 'muted')
    if idx:
        txt(slide, 12.2, 7.02, 0.5, 0.22, f'{idx:02d}', 9, False, 'muted', PP_ALIGN.RIGHT)


def make_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    rect(s, 0.7, 0.65, 11.95, 6.15, 'panel')
    txt(s, 1.25, 1.55, 10.8, 0.45, '论文课堂汇报', 15, True, 'accent', PP_ALIGN.CENTER)
    txt(s, 1.2, 2.25, 10.9, 0.95, '专利权利要求 + 引文：\n技术起源与演化分析', 31, True, 'primary', PP_ALIGN.CENTER)
    txt(s, 1.55, 3.65, 10.2, 0.55, 'Technological origination and evolution analysis by combining patent claims and citations', 13, False, 'muted', PP_ALIGN.CENTER)
    txt(s, 1.4, 5.25, 10.5, 0.35, 'Qiu & Wang · Advanced Engineering Informatics · 2023 · Surgical robot domain', 11, False, 'muted', PP_ALIGN.CENTER)


def make_agenda(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s, '汇报路线', '8–10 分钟：先讲问题，再讲方法，最后看 surgical robot case', 2)
    items = ['研究问题：专利趋势分析为什么要看“技术起源”', '核心想法：把 claim 当技术元素，再结合 citation', '方法链路：TI → TEP → Main Path → Origin Analysis', '案例结果：3313 件 surgical robot 美国授权专利', '贡献、局限与可改进方向']
    y = 1.55
    for i, it in enumerate(items, 1):
        rect(s, 1.0, y, 11.2, 0.62, 'panel')
        txt(s, 1.25, y + 0.14, 0.5, 0.2, f'{i:02d}', 12, True, 'accent')
        txt(s, 1.9, y + 0.12, 9.8, 0.24, it, 15, True, 'text')
        y += 0.82


def make_motivation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s, '研究动机：不只看趋势，还要追问“从哪来”', '传统 patent trend analysis 多看发展过程，较少挖技术 origin', 3)
    rect(s, 0.85, 1.55, 5.65, 4.65, 'panel')
    txt(s, 1.15, 1.85, 5.0, 0.35, '已有方法关注什么？', 18, True, 'primary')
    bullet(s, 1.25, 2.45, 4.85, 2.7, ['统计分析 / time-series', '文本挖掘：keyword、SAO、LDA', '引文网络：citation path、main path', '问题：多停在“发展轨迹”，缺少技术细节层面的来源解释'], 14)
    rect(s, 6.85, 1.55, 5.65, 4.65, 'panel')
    txt(s, 7.15, 1.85, 5.0, 0.35, '本文想补什么？', 18, True, 'primary')
    bullet(s, 7.25, 2.45, 4.85, 2.7, ['claim = 技术元素 / 创新点', 'citation = 技术继承关系', 'claim similarity = 引文强度', '最终揭示技术融合与演化过程'], 14)


def make_framework(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s, '整体框架：claims + citations 的四段流水线', '把专利文本的细节和引文网络的方向性接起来', 4)
    steps = [('1', 'Patent data', '收集 Title / Abstract / Claims / Citations / Issue Date'), ('2', 'TI', 'claim 相似度，得到技术继承强度'), ('3', 'TEP', '按时间窗找重要 technical elements'), ('4', 'Main Paths', '构造技术发展主路径'), ('5', 'Origin Analysis', '沿主路径追踪技术融合与来源')]
    x, y = 0.75, 2.2
    for i, (no, head, body) in enumerate(steps):
        rect(s, x, y, 2.25, 1.45, 'panel')
        txt(s, x + 0.18, y + 0.18, 1.85, 0.25, f'{no}  {head}', 12, True, 'primary')
        txt(s, x + 0.2, y + 0.65, 1.85, 0.52, body, 9, False, 'text')
        if i < 4:
            txt(s, x + 2.28, y + 0.48, 0.35, 0.3, '→', 20, True, 'accent', PP_ALIGN.CENTER)
        x += 2.48
    txt(s, 1.1, 5.25, 11.0, 0.4, '一句话：用 claim 解决“技术细节”，用 citation 解决“继承方向”。', 18, True, 'dark', PP_ALIGN.CENTER)


def make_ti(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s, '关键指标 1：Technological Inheritance (TI)', '衡量“被引专利 j 对引用专利 i 的技术继承强度”', 5)
    rect(s, 0.85, 1.45, 5.7, 4.9, 'panel')
    txt(s, 1.15, 1.78, 5.0, 0.3, '计算逻辑', 18, True, 'primary')
    bullet(s, 1.25, 2.35, 4.8, 2.8, ['把每条 claim 预处理成技术词集合', '将 claim 编码为 0–1 term vector', '对 citing patent 的每条 claim，找 cited patent 中最相似 claim', '平均最大 cosine similarity，得到 TI'], 14)
    rect(s, 6.9, 1.45, 5.55, 4.9, 'soft')
    txt(s, 7.2, 1.85, 5.0, 0.3, '公式直觉', 18, True, 'primary')
    txt(s, 7.2, 2.55, 4.85, 0.55, 'MCSV = max cosine(claim_i_m, claim_j_k)', 15, True, 'dark')
    txt(s, 7.2, 3.35, 4.85, 0.55, 'TI(i,j) = average_m MCSV(i,j,m)', 15, True, 'dark')
    txt(s, 7.2, 4.45, 4.8, 0.7, 'TI 不是简单“有没有引用”，而是问：引用关系里，到底继承了多少 claim-level 技术内容。', 13, False, 'text')


def make_tep_mainpath(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s, '关键指标 2：TEP 与 Main Path', '从技术元素重要性走向演化路径', 6)
    cols = [('TEP', ['Technical Element Persistence', '在时间区间内评估 claim 的持续影响', '用于筛选 important claims / patents']), ('RNIT', ['Relationship Network of Important Technologies', '按时间窗选 top-Q TEP patents', '再用 citation 连接重要技术']), ('Main Path', ['从 RNIT 中筛主路径', '要求路径长度与重要节点数', '最终展示技术演化轨迹'])]
    x = 0.85
    for head, items in cols:
        rect(s, x, 1.45, 3.55, 4.85, 'panel')
        txt(s, x + 0.25, 1.8, 3.0, 0.3, head, 19, True, 'primary')
        bullet(s, x + 0.3, 2.35, 2.75, 2.8, items, 13)
        x += 4.0


def make_data(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s, '案例数据：Surgical robot domain', 'USPTO 授权专利，Patsnap 检索，截止 2020-05-26', 7)
    metrics = [('3313', 'patents', '美国授权 surgical robot 专利'), ('13,097', 'citations', '领域内部 citation'), ('7.906', 'avg degree', '平均引用连接度'), ('0.122', 'median TI', 'TI 中位数'), ('0.077', 'avg TI', 'TI 平均值')]
    x = 0.85
    for num, label, desc in metrics:
        rect(s, x, 1.65, 2.25, 1.45, 'panel')
        txt(s, x + 0.12, 1.92, 2.0, 0.35, num, 23, True, 'primary', PP_ALIGN.CENTER)
        txt(s, x + 0.12, 2.42, 2.0, 0.42, f'{label}\n{desc}', 8, False, 'muted', PP_ALIGN.CENTER)
        x += 2.45
    rect(s, 1.1, 4.05, 11.1, 1.55, 'soft')
    txt(s, 1.45, 4.35, 10.4, 0.35, '解释', 17, True, 'primary')
    txt(s, 1.45, 4.88, 10.2, 0.35, '内部引用不算特别密集；TI 长尾分布说明：只有少数引用关系体现较强技术继承，多数 citation 的 claim-level 继承较弱。', 14, False, 'text')


def make_findings_elements(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s, '结果 1：重要技术元素在时间窗中迁移', '2011–2015 与 2016–2020 的 top claims 对比', 8)
    rect(s, 0.85, 1.5, 5.65, 4.75, 'panel')
    txt(s, 1.15, 1.82, 5.0, 0.3, '2011–2015', 18, True, 'primary')
    bullet(s, 1.2, 2.35, 4.85, 2.8, ['Robotic arm motion control', 'X-ray / CT imaging for surgery', 'Endoscope control', 'Bone positioning / alteration', 'System display and operator commands'], 13)
    rect(s, 6.85, 1.5, 5.65, 4.75, 'panel')
    txt(s, 7.15, 1.82, 5.0, 0.3, '2016–2020', 18, True, 'primary')
    bullet(s, 7.2, 2.35, 4.85, 2.8, ['Surgical operation systems', 'Surgical microprocessors', 'Endoscope insertion across skin', 'Minimally invasive technologies', 'Robotic arm motion control remains important'], 13)
    txt(s, 1.0, 6.45, 11.4, 0.28, 'Takeaway：热点从机械/影像控制，逐步扩展到系统化、微处理器与微创操作。', 13, True, 'accent', PP_ALIGN.CENTER)


def make_paths(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s, '结果 2：四条 surgical robot 技术主路径', '主路径把复杂 citation network 压缩成可讲的演化路线', 9)
    paths = [('Path 1', '定位与切割', 'Laser surgical knife → Robot-aided surgery → Surgical robotic tools'), ('Path 2', 'X-ray / scanning', 'Stereotactic apparatus → CT-assisted robotic system → real-time navigational guidance'), ('Path 3', '眼科 / 头部手术', 'Fluid control → ophthalmic suction → probe position + head imaging'), ('Path 4', '骨科定位操作', 'Bone alteration → image-directed robotic surgery → navigation guidance')]
    y = 1.45
    for pname, tag, body in paths:
        rect(s, 0.95, y, 11.45, 0.9, 'panel')
        txt(s, 1.2, y + 0.22, 1.0, 0.25, pname, 12, True, 'accent')
        txt(s, 2.25, y + 0.22, 1.6, 0.25, tag, 14, True, 'primary')
        txt(s, 4.05, y + 0.18, 7.8, 0.3, body, 11, False, 'text')
        y += 1.08


def make_fusion_conclusion(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s, '结论：claim-level 视角让“技术融合”可解释', '贡献、发现与局限', 10)
    cols = [('理论贡献', ['TI：claim-based inheritance', 'TEP：technical element persistence', 'Main path：基于技术细节的演化路径', 'Origin analysis：解释融合来源']), ('案例发现', ['未来热点：数字终端、内窥镜、微创机器人', '四条主路径：切割、扫描、头部、骨科', '融合主线：定位 + 影像 + 导航']), ('局限', ['未考虑 claim 结构', 'citation 有时间滞后', 'cosine similarity 偏浅，可换更强 NLP', '数据集还可扩展'])]
    x = 0.85
    for head, items in cols:
        rect(s, x, 1.45, 3.65, 4.8, 'panel')
        txt(s, x + 0.25, 1.78, 3.1, 0.3, head, 17, True, 'primary')
        bullet(s, x + 0.3, 2.3, 2.9, 2.8, items, 12)
        x += 4.0


def build():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    make_title(prs)
    make_agenda(prs)
    make_motivation(prs)
    make_framework(prs)
    make_ti(prs)
    make_tep_mainpath(prs)
    make_data(prs)
    make_findings_elements(prs)
    make_paths(prs)
    make_fusion_conclusion(prs)
    prs.save(OUT)
    print(f'Wrote {OUT}')


if __name__ == '__main__':
    build()
