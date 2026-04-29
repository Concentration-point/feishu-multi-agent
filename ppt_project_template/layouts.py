from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from theme import SLIDE_W, SLIDE_H, MARGIN_X, MARGIN_Y, TITLE_SIZE, SUBTITLE_SIZE, BODY_SIZE, SMALL_SIZE, rgb

BLANK = 6

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_textbox(slide, x, y, w, h, text, theme, size=BODY_SIZE, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.name = theme.font_head if bold else theme.font_body
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = rgb(color or theme.text)
    return box


def add_title(slide, title, theme, subtitle=None):
    add_textbox(slide, MARGIN_X, Inches(0.45), Inches(10.5), Inches(0.55), title, theme, TITLE_SIZE, True, theme.primary)
    if subtitle:
        add_textbox(slide, MARGIN_X, Inches(1.05), Inches(10.8), Inches(0.35), subtitle, theme, SUBTITLE_SIZE, False, theme.muted)


def add_footer(slide, theme, idx=None):
    txt = f"{idx:02d}" if idx is not None else ""
    add_textbox(slide, Inches(12.15), Inches(7.0), Inches(0.55), Inches(0.22), txt, theme, SMALL_SIZE, False, theme.muted, PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, theme, fill=None, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill or theme.panel)
    shape.line.color.rgb = rgb(line or theme.panel)
    return shape


def add_bullet_list(slide, x, y, w, h, items, theme, size=BODY_SIZE):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(item)
        p.level = 0
        p.font.name = theme.font_body
        p.font.size = size
        p.font.color.rgb = rgb(theme.text)
    return box
