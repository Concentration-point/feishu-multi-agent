from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont
import json, zipfile, re, textwrap, math, os

OUT = Path('deckforge_output/surgical_robot_patent_evolution_v2')
IMG = OUT / 'images'
OUT.mkdir(parents=True, exist_ok=True)
IMG.mkdir(parents=True, exist_ok=True)

# ---------- Style Bible / spec_lock ----------
W, H = 13.333, 7.5
NAVY = RGBColor(10, 20, 36)
NAVY2 = RGBColor(18, 34, 58)
BG = RGBColor(247, 250, 253)
CARD = RGBColor(255, 255, 255)
CYAN = RGBColor(27, 183, 224)
BLUE = RGBColor(45, 104, 214)
ORANGE = RGBColor(241, 142, 43)
GREEN = RGBColor(43, 166, 117)
PURPLE = RGBColor(132, 87, 214)
RED = RGBColor(213, 78, 74)
TEXT = RGBColor(29, 40, 56)
MUTED = RGBColor(92, 108, 128)
GRID = RGBColor(223, 232, 242)
WHITE = RGBColor(255, 255, 255)

FONT = 'Microsoft YaHei'
FONT_EN = 'Arial'
FONT_MONO = 'Consolas'

spec_lock = {
    'format': 'ppt169',
    'colors': {'primary':'#0A1424','secondary':'#12223A','background':'#F7FAFD','text':'#1D2838','muted':'#5C6C80','accent_cyan':'#1BB7E0','accent_orange':'#F18E2B','accent_green':'#2BA675','accent_purple':'#8457D6'},
    'typography': {'title_family':'Microsoft YaHei','body_family':'Microsoft YaHei','mono_family':'Consolas','body_size':15},
    'page_rhythm': {'P01':'anchor','P02':'dense','P03':'dense','P04':'dense','P05':'dense','P06':'dense','P07':'dense','P08':'dense','P09':'dense','P10':'dense','P11':'dense','P12':'dense','P13':'breathing','P14':'dense','P15':'anchor'},
    'images': {'style_anchor':'clean academic technology presentation, white/deep navy, cyan/orange accents, vector-like, no text in generated assets'},
    'quality_rules': ['no raw complex formulas in PPT text boxes', 'formula pages use PNG formula panels', 'every content page has at least two modules and one visual anchor']
}
(OUT/'spec_lock.json').write_text(json.dumps(spec_lock, ensure_ascii=False, indent=2), encoding='utf-8')

style_bible = {
    'name':'DeckForge Academic Tech v2',
    'scenario':'15-minute classroom paper report',
    'layout_density':'dense by default; breathing only for synthesis pages',
    'color_use':'60% light background, 30% navy structure, 10% cyan/orange accents',
    'formula_policy':'all formulas rendered as images to avoid PPT character corruption',
    'image_policy':'no standalone AI images generated for this rebuild; vector diagrams + formula PNGs used for controllability'
}
(OUT/'style-bible.json').write_text(json.dumps(style_bible, ensure_ascii=False, indent=2), encoding='utf-8')

# ---------- Assets ----------
def font_path(name):
    candidates = [
        r'C:\Windows\Fonts\msyh.ttc', r'C:\Windows\Fonts\arial.ttf', r'C:\Windows\Fonts\consola.ttf',
        '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf'
    ]
    for p in candidates:
        if Path(p).exists(): return p
    return None

BASE_FONT = font_path('base')
MONO_FONT = r'C:\Windows\Fonts\consola.ttf' if Path(r'C:\Windows\Fonts\consola.ttf').exists() else BASE_FONT

def make_formula_image(filename, title, lines, notes=None, accent=(27,183,224)):
    scale = 2
    w, h = 1600, 780
    im = Image.new('RGB', (w, h), (255,255,255))
    d = ImageDraw.Draw(im)
    f_title = ImageFont.truetype(BASE_FONT, 48) if BASE_FONT else ImageFont.load_default()
    f_formula = ImageFont.truetype(MONO_FONT, 42) if MONO_FONT else ImageFont.load_default()
    f_note = ImageFont.truetype(BASE_FONT, 28) if BASE_FONT else ImageFont.load_default()
    # border and header
    d.rounded_rectangle([20,20,w-20,h-20], radius=30, outline=(218,230,242), width=4, fill=(255,255,255))
    d.rectangle([20,20,w-20,112], fill=(10,20,36))
    d.rectangle([20,112,260,122], fill=accent)
    d.text((60,48), title, fill=(255,255,255), font=f_title)
    y = 165
    for line in lines:
        d.rounded_rectangle([65,y-10,w-65,y+58], radius=12, fill=(246,250,253), outline=(222,232,242), width=2)
        d.text((90,y), line, fill=(20,32,48), font=f_formula)
        y += 95
    if notes:
        y += 10
        for note in notes:
            d.ellipse([78,y+9,92,y+23], fill=accent)
            d.text((110,y), note, fill=(65,82,104), font=f_note)
            y += 45
    path = IMG / filename
    im.save(path)
    return str(path)

formula_ti = make_formula_image('formula_ti.png', 'TI: Technological Inheritance', [
    'cos(key_i_m, key_j_k) = dot(key_i_m, key_j_k) / (norm_i_m * norm_j_k)',
    'MCSV(i,j,m) = max_k cos(key_i_m, key_j_k)',
    'TI(i,j) = (1 / M_i) * sum_m MCSV(i,j,m)'
], ['For each claim in citing patent i, find the most similar claim in cited patent j.', 'TI is the average maximum claim similarity on a citation link.'])
formula_tep = make_formula_image('formula_tep.png', 'TEP: Technical Element Persistence', [
    'TICV(i,j,m,l) = MCSV(i,j,m)',
    'TEP(p,u,alpha) = sum over backward citation paths of TICV(path,s,u)',
    'TEP(p,alpha) = sum_m TEP(p,m,alpha)'
], ['A claim is important if later patents repeatedly inherit it.', 'Indirect citation contribution propagates along the citation path.'], accent=(241,142,43))
formula_path = make_formula_image('formula_path.png', 'Main Path Construction', [
    'PT_Q(alpha) = top Q patents ranked by TEP in interval alpha',
    'PT_Q = union over alpha of PT_Q(alpha)',
    'path is kept if length >= a and important_patents >= b'
], ['Case parameters: Q = 15, a = 4, b = 3.', 'RNIT connects important technologies through citation paths.'], accent=(132,87,214))

# ---------- PPT helpers ----------
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W), Inches(H)

slide_specs=[]

def rgb(hexstr):
    hexstr=hexstr.lstrip('#')
    return RGBColor(int(hexstr[0:2],16), int(hexstr[2:4],16), int(hexstr[4:6],16))

def add_bg(slide, color=BG):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color


def set_run(run, size=15, color=TEXT, bold=False, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def text_box(slide, text, x,y,w,h, size=15, color=TEXT, bold=False, align=None, font=FONT, auto=True):
    box=slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf=box.text_frame; tf.clear(); tf.word_wrap=True
    if auto: tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    lines=text.split('\n')
    for i,line in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        if align: p.alignment=align
        p.space_after=Pt(2)
        r=p.add_run(); r.text=line; set_run(r,size,color,bold,font)
    return box


def rect(slide,x,y,w,h,fill=CARD,line=GRID,radius=True):
    shp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x),Inches(y),Inches(w),Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb=fill
    shp.line.color.rgb=line; shp.line.width=Pt(0.8)
    return shp


def title(slide, t, takeaway, page):
    rect(slide,0,0,13.333,0.72,NAVY,NAVY,False)
    text_box(slide,t,0.48,0.17,8.4,0.35,18,RGBColor(255,255,255),True)
    text_box(slide,takeaway,8.35,0.18,4.35,0.34,9.5,RGBColor(188,213,232),False,PP_ALIGN.RIGHT)
    # footer
    text_box(slide,'DeckForge v2 · Patent claims + citations · Surgical robot domain',0.48,7.18,8.0,0.18,7.8,MUTED)
    text_box(slide,f'{page:02d}/15',12.18,7.16,0.68,0.2,8.5,MUTED,True,PP_ALIGN.RIGHT)


def card(slide, title_, body, x,y,w,h, accent=CYAN, title_size=12.5, body_size=10.5):
    rect(slide,x,y,w,h,CARD,GRID,True)
    bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x),Inches(y),Inches(0.07),Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb=accent; bar.line.fill.background()
    text_box(slide,title_,x+0.18,y+0.13,w-0.3,0.25,title_size,NAVY,True)
    text_box(slide,body,x+0.18,y+0.48,w-0.3,h-0.56,body_size,TEXT)


def tag(slide,text,x,y,color=CYAN,w=1.15):
    shp=rect(slide,x,y,w,0.28,rgb('#EAF7FC'),color,True)
    shp.line.width=Pt(0.4)
    text_box(slide,text,x+0.07,y+0.06,w-0.14,0.13,7.5,color,True,PP_ALIGN.CENTER)


def bullet_list(slide, items, x,y,w,h, size=11.5, accent=CYAN):
    yy=y
    for item in items:
        dot=slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(yy+0.07), Inches(0.09), Inches(0.09))
        dot.fill.solid(); dot.fill.fore_color.rgb=accent; dot.line.fill.background()
        text_box(slide,item,x+0.18,yy,w-0.18,0.25,size,TEXT)
        yy += 0.34


def arrow(slide,x1,y1,x2,y2,color=CYAN,width=1.4):
    c=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    c.line.color.rgb=color; c.line.width=Pt(width); c.line.end_arrowhead=True
    return c


def node(slide,label,x,y,w=1.2,h=0.46,fill=CARD,line=CYAN,size=9.2):
    rect(slide,x,y,w,h,fill,line,True)
    text_box(slide,label,x+0.05,y+0.11,w-0.1,h-0.12,size,NAVY,True,PP_ALIGN.CENTER)


def record(page,title_,takeaway,rhythm='dense'):
    slide_specs.append({'page':page,'title':title_,'takeaway':takeaway,'rhythm':rhythm})

# ---------- Slides ----------
# 1 cover
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s,NAVY)
# dense visual network
for i in range(9):
    x=7.3+(i%3)*1.55; y=1.35+(i//3)*1.25
    node(s, f'P{i+1}', x,y,0.65,0.38, RGBColor(236,248,252) if i in [0,2,4,6] else RGBColor(255,247,237), CYAN if i in [0,2,4,6] else ORANGE, 8)
for a,b in [((7.95,1.54),(8.85,1.54)),((9.5,1.54),(10.4,1.54)),((8.0,1.75),(8.8,2.78)),((9.5,1.75),(10.4,2.78)),((7.95,2.8),(8.85,4.0)),((9.5,2.8),(10.4,4.0)),((8.0,4.02),(10.4,4.02))]: arrow(s,*a,*b,RGBColor(79,108,138),1.1)
text_box(s,'基于 Patent Claims 与 Citations 的\n技术起源与演化分析',0.65,1.1,6.8,1.35,28,RGBColor(255,255,255),True)
text_box(s,'以手术机器人领域为例｜Advanced Engineering Informatics 58 (2023) 102145',0.68,2.7,6.55,0.28,12,CYAN,False)
card(s,'汇报定位','15 分钟课堂汇报\n重点讲清：问题、方法、案例发现、评价',0.7,4.25,2.7,1.25,CYAN)
card(s,'论文核心','Claims = technical elements\nCitations = inheritance links',3.65,4.25,2.95,1.25,ORANGE)
card(s,'方法关键词','TI · TEP · RNIT · Main Paths\nOrigin & Development Analysis',0.7,5.75,5.9,0.82,PURPLE)
text_box(s,'01/15',12.25,7.05,0.7,0.2,8.5,RGBColor(180,204,225),True,PP_ALIGN.RIGHT)
record(1,'封面','用 claim-level 视角重新解释技术演化','anchor')

# 2 brief
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'Deck Brief：15 分钟讲什么？','目标：让听众记住“技术起源分析”比普通趋势分析更细',2)
card(s,'一句话主旨','这篇论文把专利分析从“专利引用网络”推进到“权利要求技术元素网络”。',0.65,1.15,4.05,1.0,CYAN,13,11)
card(s,'课堂听众需要带走','1. Claims 为什么适合当技术元素\n2. TI/TEP 如何工作\n3. 手术机器人领域发现了什么\n4. 方法有什么短板',4.95,1.15,3.55,1.75,ORANGE,13,10.5)
card(s,'15 页节奏','背景动机 3 页\n方法框架 6 页\n案例发现 4 页\n评价总结 2 页',8.75,1.15,3.8,1.75,PURPLE,13,10.5)
# timeline
for i,(lab,col) in enumerate([('Why',CYAN),('How',ORANGE),('Case',GREEN),('So what',PURPLE)]):
    x=1.0+i*3.0
    node(s,lab,x,4.25,1.15,0.46,RGBColor(240,248,252),col,10)
    if i<3: arrow(s,x+1.15,4.48,x+3.0,4.48,col,1.2)
text_box(s,'页 2–4：研究缺口和框架',0.8,5.0,2.45,0.3,10,MUTED,False,PP_ALIGN.CENTER)
text_box(s,'页 5–10：核心指标和路径构建',3.7,5.0,2.7,0.3,10,MUTED,False,PP_ALIGN.CENTER)
text_box(s,'页 11–13：手术机器人案例',6.88,5.0,2.45,0.3,10,MUTED,False,PP_ALIGN.CENTER)
text_box(s,'页 14–15：评价与结论',9.73,5.0,2.45,0.3,10,MUTED,False,PP_ALIGN.CENTER)
card(s,'Takeaway','不要把这篇论文讲成“公式合集”；它真正的贡献是让主路径分析能解释技术从哪里来。',1.1,5.85,11.1,0.72,GREEN,12,11)
record(2,'Deck Brief','围绕 why/how/case/so what 四段讲','dense')

# 3 gap
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'研究缺口：趋势分析很多，起源分析不足','已有方法能看“怎么发展”，但难解释“从哪里来”',3)
card(s,'Common Trend Analysis','统计分析、时间序列、机器学习、Delphi。\n优势：宏观趋势清楚。\n短板：技术细节弱。',0.65,1.2,3.2,2.0,BLUE)
card(s,'Patent Text Mining','关键词、LDA、SAO、聚类、文本网络。\n优势：能发现技术热点。\n短板：继承关系弱。',4.05,1.2,3.2,2.0,CYAN)
card(s,'Citation Main Path','引用网络、主路径、知识流动。\n优势：演化路径明确。\n短板：不知道引用了什么。',7.45,1.2,3.2,2.0,ORANGE)
card(s,'本文补位','用 claims 捕捉技术元素，\n用 citations 约束继承关系，\n再追踪主路径上的技术融合。',10.85,1.2,1.85,2.0,PURPLE,11,9.2)
# gap funnel
text_box(s,'普通趋势分析',1.0,4.05,2.0,0.3,13,NAVY,True,PP_ALIGN.CENTER); arrow(s,3.0,4.2,4.5,4.2,MUTED,1.1)
text_box(s,'专利引用路径',4.65,4.05,2.0,0.3,13,NAVY,True,PP_ALIGN.CENTER); arrow(s,6.65,4.2,8.15,4.2,MUTED,1.1)
text_box(s,'Claim-level 起源',8.3,4.05,2.3,0.3,13,NAVY,True,PP_ALIGN.CENTER)
card(s,'研究问题','如何结合 patent claims 和 citations，揭示一个技术领域的 technological origination、technology fusion 和 evolution process？',1.2,5.35,10.95,0.9,GREEN,13,11.2)
record(3,'研究缺口','论文补足“技术起源与融合”解释','dense')

# 4 framework
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'总体框架：Claims + Citations 的五段流水线','对应论文 Fig. 1，重画为课堂版',4)
steps=[('1 Patent\nCollection','Patsnap\nUSPTO',0.55,CYAN),('2 Claims\nPreprocess','0-1 technical\nterm matrix',2.65,BLUE),('3 TI\nInheritance','claim similarity\non citation',4.75,ORANGE),('4 TEP\nImportance','persistence in\ntime interval',6.85,PURPLE),('5 Main Paths\n& Origins','RNIT + origin\ndevelopment',8.95,GREEN)]
for title_,sub,x,col in steps:
    node(s,title_,x,1.4,1.52,0.7,RGBColor(255,255,255),col,9.5)
    text_box(s,sub,x,2.2,1.52,0.4,8.5,MUTED,False,PP_ALIGN.CENTER)
for i in range(len(steps)-1): arrow(s,steps[i][2]+1.52,1.76,steps[i+1][2],1.76,CYAN,1.25)
card(s,'输入层','3313 件手术机器人 USPTO 授权专利；字段包括 title、abstract、claims、citations、issue date。',0.75,3.25,3.75,1.1,CYAN)
card(s,'计算层','把每个 claim 变成技术词向量；在引用关系上计算 claim-to-claim similarity。',4.8,3.25,3.75,1.1,ORANGE)
card(s,'解释层','用 TEP 找重要技术元素，用 RNIT 与 main paths 展示演化，再追踪 claim-level origin。',8.85,3.25,3.75,1.1,PURPLE)
card(s,'方法定位','不是“更复杂的引用网络”，而是在引用网络里嵌入技术细节。',1.35,5.3,10.55,0.78,GREEN,13,11.2)
record(4,'总体框架','五段流程：收集、预处理、TI、TEP、主路径/起源','dense')

# 5 claim representation
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'Claims 如何变成 technical elements？','每个 claim 被当作一个可追踪的微观技术元素',5)
card(s,'为什么选 claims','Claims 是专利保护范围，最接近“创新点”；比 abstract 更细，比 description 更聚焦。',0.65,1.1,3.7,1.1,CYAN)
card(s,'预处理','去除 the / wherein / mean 等低信息词；提取技术词，构建 Dict。',0.65,2.45,3.7,1.1,ORANGE)
card(s,'矩阵化','每个 claim 转成 0-1 向量：词出现为 1，不出现为 0；同词多次出现只算一次。',0.65,3.8,3.7,1.25,PURPLE)
# matrix visual
x0,y0=5.0,1.25
text_box(s,'Claim × Technical Terms Matrix',x0,1.0,5.9,0.28,13,NAVY,True)
terms=['position','robot','arm','bone','image','control','endoscope']
for i,t in enumerate(terms): text_box(s,t,x0+1.25+i*0.72,y0,0.62,0.3,7.5,MUTED,True,PP_ALIGN.CENTER)
claims=['Claim 1','Claim 2','Claim 3','Claim 4']
vals=[[1,1,0,1,0,0,0],[0,1,1,0,0,1,0],[0,0,0,0,1,1,1],[1,0,0,1,1,0,0]]
for r,cname in enumerate(claims):
    text_box(s,cname,x0,y0+0.55+r*0.58,1.0,0.25,8.5,NAVY,True)
    for c,v in enumerate(vals[r]):
        fill=CYAN if v else RGBColor(235,241,247)
        shp=slide= s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x0+1.25+c*0.72), Inches(y0+0.52+r*0.58), Inches(0.38), Inches(0.28))
        shp.fill.solid(); shp.fill.fore_color.rgb=fill; shp.line.color.rgb=RGBColor(220,230,240)
        text_box(s,str(v),x0+1.25+c*0.72,y0+0.56+r*0.58,0.38,0.13,7,WHITE if v else MUTED,True,PP_ALIGN.CENTER)
card(s,'课堂讲法','这一步不用讲成 NLP 细节；告诉听众：作者把 claim 变成“技术词是否出现”的向量，然后比较相似度。',5.05,5.35,6.75,0.82,GREEN,12.5,10.5)
record(5,'Claims 表示','claims 被转成技术词 0-1 矩阵','dense')

# 6 TI
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'TI：引用关系中的技术继承强度','用平均最大 claim 相似度判断“引用到底继承了多少”',6)
s.shapes.add_picture(formula_ti, Inches(0.65), Inches(1.1), Inches(6.4), Inches(3.1))
card(s,'读公式的方法','1. 对引用专利 i 的每个 claim。\n2. 在被引专利 j 的 claims 里找最像的一个。\n3. 把这些最大相似值平均，得到 TI(i,j)。',7.35,1.1,5.25,1.55,CYAN)
card(s,'为什么重要','普通 citation 只知道“谁引用谁”；TI 进一步回答“具体继承强不强”。',7.35,2.9,5.25,1.0,ORANGE)
card(s,'论文处理','TI > 0.7 的引用被剔除，因为可能是同一技术连续更新，novelty 较少。',7.35,4.15,5.25,0.95,PURPLE)
# mini claim link
for i,l in enumerate(['j1','j2','j3']): node(s,l,1.05,4.85+i*0.42,0.5,0.28,RGBColor(239,246,255),BLUE,7.5)
for i,l in enumerate(['i1','i2','i3']): node(s,l,3.3,4.85+i*0.42,0.5,0.28,RGBColor(255,247,237),ORANGE,7.5)
arrow(s,1.55,4.99,3.3,5.39,CYAN,1.2); arrow(s,1.55,5.41,3.3,4.99,PURPLE,1.0); arrow(s,1.55,5.82,3.3,5.82,GREEN,1.0)
text_box(s,'claim-to-claim max matching',1.45,6.15,2.6,0.2,8.5,MUTED,False,PP_ALIGN.CENTER)
record(6,'TI 指标','TI 是引用边上的技术继承强度','dense')

# 7 origin
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'Main Technical Origin：每个 claim 的来源从哪里来？','MSCI 找到一个 claim 在所有被引专利中的主技术起源',7)
card(s,'问题','一个专利通常引用多个专利；某个 claim 的真正来源可能藏在其中某一个被引专利里。',0.65,1.15,3.75,1.1,CYAN)
card(s,'方法','对 claim i_m 遍历所有 cited patents，取全局最大 MCSV，对应 claim 即 MSCI。',0.65,2.55,3.75,1.1,ORANGE)
card(s,'价值','研发人员可以看到新技术由哪些旧技术元素组合而成，也能找到可替代技术来源。',0.65,3.95,3.75,1.1,GREEN)
# visual fan-in
node(s,'Claim i_m\nnew element',9.4,3.0,1.5,0.62,RGBColor(255,247,237),ORANGE,9)
for k,(lab,y,col) in enumerate([('Patent A\nclaim a2',1.35,BLUE),('Patent B\nclaim b1',2.35,CYAN),('Patent C\nclaim c4',4.05,PURPLE),('Patent D\nclaim d3',5.05,GREEN)]):
    node(s,lab,5.45,y,1.35,0.58,WHITE,col,8.5)
    arrow(s,6.8,y+0.29,9.4,3.31,col,1.0+0.18*k)
tag(s,'max similarity',7.65,2.75,ORANGE,1.35)
card(s,'案例口径','论文用 US5887121 展示：机器人关节控制 claim 可追溯到 US5297057 / US5377310 等早期技术。',5.35,5.85,6.7,0.7,PURPLE,12,10)
record(7,'Claim 主技术起源','MSCI 把 claim 追溯到最相似的被引 claim','dense')

# 8 TEP
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'TEP：哪个技术元素在某时期最重要？','一个 claim 被后续技术持续继承，TEP 就越高',8)
s.shapes.add_picture(formula_tep, Inches(0.65), Inches(1.1), Inches(6.35), Inches(3.1))
card(s,'直观解释','如果某个旧 claim 不断被后续专利直接或间接继承，它在该时间段内就是重要技术元素。',7.3,1.1,5.25,1.05,CYAN)
card(s,'时间区间','论文按 5 年划分时间段，如 2011–2015、2016–2020。每个区间单独计算重要元素。',7.3,2.4,5.25,1.05,ORANGE)
card(s,'贡献传播','直接引用贡献为 MCSV；间接引用沿路径传播，体现技术影响力随路径传递。',7.3,3.7,5.25,1.05,PURPLE)
# propagation chain
for i,(lab,col) in enumerate([('old\nclaim',BLUE),('middle\nclaim',CYAN),('later\nclaim',ORANGE),('TEP\nalpha',GREEN)]):
    node(s,lab,1.0+i*1.35,5.25,0.9,0.5,WHITE,col,8)
    if i<3: arrow(s,1.9+i*1.35,5.5,2.35+i*1.35,5.5,col,1.1)
text_box(s,'TICV 累计 + 间接传播',1.6,6.05,3.6,0.25,9.5,MUTED,False,PP_ALIGN.CENTER)
record(8,'TEP 指标','TEP 用被继承程度衡量 claim 重要性','dense')

# 9 main paths
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'Main Paths：从重要技术元素到演化主路径','先选重要专利，再构建 RNIT，而不是全网乱画',9)
s.shapes.add_picture(formula_path, Inches(0.65), Inches(1.1), Inches(6.25), Inches(2.95))
card(s,'实验参数','Q = 15：每个 5 年区间选 TEP 前 15 专利。\na = 4：路径长度至少 4。\nb = 3：路径中重要专利至少 3。',7.15,1.1,5.35,1.4,CYAN)
card(s,'RNIT','Relationship Network of Important Technologies：重要技术关系网络，由重要专利和连接它们的引用路径组成。',7.15,2.8,5.35,1.15,ORANGE)
card(s,'筛选原则','如果多条路径满足条件，比较路径上专利的 forward citation numbers，保留影响力更强的路径。',7.15,4.25,5.35,1.05,PURPLE)
# mini RNIT graph
coords=[(1.0,5.55,'P1',ORANGE),(2.2,5.1,'P2',WHITE),(3.4,5.55,'P3',ORANGE),(4.6,5.1,'P4',ORANGE),(5.8,5.55,'P5',WHITE)]
for x,y,l,c in coords: node(s,l,x,y,0.5,0.3,c,CYAN if c==WHITE else ORANGE,7)
for i in range(len(coords)-1): arrow(s,coords[i][0]+0.5,coords[i][1]+0.15,coords[i+1][0],coords[i+1][1]+0.15,MUTED,1)
tag(s,'orange = important patent',2.4,6.2,ORANGE,1.85)
record(9,'主路径构建','用 TEP 选重要专利，再筛主路径','dense')

# 10 origin development
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'Origin & Development Analysis：把主路径拆到 claim-level','主路径回答“哪条路”，起源分析回答“由哪些技术融合”',10)
# central path
for i,(lab,col) in enumerate([('Patent A',BLUE),('Patent B',CYAN),('Patent C',ORANGE),('Patent D',PURPLE)]):
    node(s,lab,1.05+i*2.8,2.05,1.1,0.46,WHITE,col,9)
    if i<3: arrow(s,2.15+i*2.8,2.28,3.85+i*2.8,2.28,col,1.25)
# claims beneath
for i,x in enumerate([1.05,3.85,6.65,9.45]):
    for j in range(3):
        fill=RGBColor(238,248,252) if j==0 else RGBColor(248,250,253)
        node(s,f'c{j+1}',x+0.1+j*0.38,3.05,0.3,0.24,fill,CYAN if j==0 else GRID,6)
card(s,'步骤 1','取主路径上的重要专利 IP_path。',0.75,4.35,2.6,0.75,CYAN,12,10)
card(s,'步骤 2','对每个重要专利选 TEP 前 c 个 claims；案例中 c = 5。',3.55,4.35,3.0,0.75,ORANGE,12,10)
card(s,'步骤 3','为这些 claims 找 MSCI；必要时补相邻专利之间最相似 claim 对。',6.75,4.35,3.2,0.75,PURPLE,12,10)
card(s,'输出','得到一张“主路径 + 外部技术来源”的融合图。',10.15,4.35,2.45,0.75,GREEN,12,10)
card(s,'这一步的课堂价值','它把专利路径讲成了技术故事：某个机器人系统不是凭空出现，而是定位、影像、控制、导航等元素持续汇入。',1.0,5.85,11.4,0.72,GREEN,12.5,10.6)
record(10,'起源与发展分析','从专利级路径深入到 claim 级技术融合','dense')

# 11 data results
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'案例数据：Surgical Robot Domain','3313 件 USPTO 授权专利，13097 条内部引用',11)
metrics=[('3313','USPTO 授权专利',CYAN),('13097','内部引用',ORANGE),('7.906','Average degree',PURPLE),('0.122','Median TI',GREEN),('0.077','Average TI',BLUE)]
for i,(num,lab,col) in enumerate(metrics):
    x=0.65+i*2.52
    rect(s,x,1.1,2.15,1.0,CARD,col,True)
    text_box(s,num,x+0.08,1.24,1.95,0.35,20,col,True,PP_ALIGN.CENTER)
    text_box(s,lab,x+0.08,1.68,1.95,0.22,8.8,MUTED,False,PP_ALIGN.CENTER)
# histogram schematic bigger
text_box(s,'TI Distribution：long-tail',0.9,2.8,3.2,0.28,14,NAVY,True)
base_x,base_y=1.0,5.75
bars=[1.65,2.05,1.78,1.25,0.85,0.55,0.36,0.22,0.14,0.09]
for i,h in enumerate(bars):
    shp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(base_x+i*0.36), Inches(base_y-h), Inches(0.22), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb=CYAN if i<4 else RGBColor(178,196,214); shp.line.fill.background()
text_box(s,'多数 citation 技术继承度较低；少数高 TI 引用构成长尾。',0.9,6.05,4.4,0.25,10.5,MUTED)
card(s,'解释 1','很多引用是背景性或法律性引用，不一定意味着强技术继承。',5.6,2.75,3.2,1.0,CYAN)
card(s,'解释 2','高 TI 引用较少，也可能因为高度相似专利存在侵权/novelty 风险。',9.1,2.75,3.2,1.0,ORANGE)
card(s,'汇报点','这正是本文要结合 claims 的原因：不能把所有 citation 都当成同等强度的技术继承。',5.6,4.2,6.7,1.0,PURPLE)
record(11,'案例数据','TI 长尾说明 citation 强度需要细分','dense')

# 12 elements trend
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'重要技术元素变化：2011–2015 vs 2016–2020','TEP 排名揭示手术机器人热点迁移',12)
card(s,'2011–2015｜偏机械控制与影像','• robotic arm motion control\n• X-ray imaging / CT scanner\n• endoscope control\n• bone alteration positioning\n• surgical trajectory movement',0.65,1.15,5.7,2.25,BLUE,13,10.5)
card(s,'2016–2020｜偏系统集成与微创器械','• surgical operation systems\n• microprocessor control device\n• endoscope insertion with trocar sleeves\n• minimally invasive surgical robot\n• stapling device technique',6.95,1.15,5.7,2.25,ORANGE,13,10.5)
# rank movement
text_box(s,'技术热度变化信号',0.85,4.05,2.2,0.25,13,NAVY,True)
for i,(lab,lrank,rrank,col) in enumerate([('US5236432: device display & operator command',5,2,CYAN),('US4979949: robotic arm movement & positioning',4,8,ORANGE),('US5217003: endoscope with trocar sleeves',10,3,GREEN)]):
    y=4.55+i*0.55
    text_box(s,lab,1.0,y,4.2,0.22,9.5,TEXT)
    node(s,str(lrank),5.65,y-0.02,0.35,0.24,WHITE,col,7)
    arrow(s,6.0,y+0.1,7.5,y+0.1,col,1)
    node(s,str(rrank),7.55,y-0.02,0.35,0.24,WHITE,col,7)
text_box(s,'2011–2015 rank',5.35,4.18,1.0,0.2,8,MUTED,False,PP_ALIGN.CENTER)
text_box(s,'2016–2020 rank',7.2,4.18,1.1,0.2,8,MUTED,False,PP_ALIGN.CENTER)
card(s,'可能机会','microprocessor 与 stapling technique 在后期进入重要技术元素列表，论文认为值得关注。',8.9,4.4,3.45,1.2,PURPLE,12,10.5)
record(12,'重要技术元素','热点从机械/影像控制转向系统集成、微处理器和微创器械','dense')

# 13 four paths breathing but with visual
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'四条主要技术演化路径','Table 5 的课堂化重构：四条路就是四个子领域故事',13)
paths=[('Path 1','定位与切割','Laser knife → Laser surgery → Robot-aided surgery → robotic tools/data architecture',RED),('Path 2','扫描 / 影像导航','Stereotactic apparatus → CT-assisted surgery → image-directed robot → navigation guidance',BLUE),('Path 3','眼科 / 头部手术','Fluid control → intraocular suction → ophthalmic instruments → head probe positioning',GREEN),('Path 4','骨科手术','Bone alteration → image-directed robot → computer-assisted skeletal surgery → fluoroscopic navigation',PURPLE)]
for i,(p,t,b,col) in enumerate(paths):
    y=1.25+i*1.28
    rect(s,0.85,y,11.6,0.88,CARD,col,True)
    text_box(s,p,1.05,y+0.18,0.95,0.22,11,col,True)
    text_box(s,t,2.05,y+0.13,1.55,0.28,13,NAVY,True)
    text_box(s,b,3.85,y+0.16,7.95,0.25,10.5,TEXT)
    tag(s,'main path',11.35,y+0.3,col,0.85)
card(s,'汇报抓手','这页不要逐个专利死背；按“定位切割、影像导航、眼科头部、骨科”四个应用方向讲。',1.35,6.35,10.7,0.58,CYAN,12,10)
record(13,'四条主路径','四条路径对应四类手术机器人技术故事','breathing')

# 14 path2
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'重点深挖：Path 2 的技术融合过程','扫描/影像导航路径最能体现 claim-level origin analysis 的价值',14)
main=[('US4341220\nStereotactic',0.65),('US4791934\nCT-assisted',3.05),('US5086401\nImage-directed',5.45),('US5682886\nComputer-assisted',7.95),('US6470207\nNavigation',10.45)]
for lab,x in main: node(s,lab,x,3.05,1.45,0.66,WHITE,BLUE,8.2)
for i in range(len(main)-1): arrow(s,main[i][1]+1.45,3.38,main[i+1][1],3.38,BLUE,1.35)
origins=[('positioning\ntechnique',0.75,1.45,0),('X-ray\nscanning',1.65,4.75,0),('picture\nfusion',3.35,1.45,1),('bone\npositioning',5.7,4.75,2),('volumetric\nmodel',7.8,1.45,3),('anatomical\nfeatures',9.55,4.75,4),('virtual\nguidewire',11.25,1.45,4)]
for lab,x,y,target in origins:
    node(s,lab,x,y,1.05,0.5,RGBColor(239,247,253),CYAN,7.4)
    tx=main[target][1]+0.72; ty=3.05 if y<3 else 3.71
    arrow(s,x+0.52,y+0.5 if y<3 else y,tx,ty,CYAN,0.9)
card(s,'演化逻辑','从 stereotactic + X-ray scanning 起步，逐步融合 picture fusion、digital imaging、bone positioning、volumetric model、3D imaging 和 navigational guidance。',0.95,5.85,11.35,0.88,ORANGE,12,10.3)
record(14,'Path 2 深挖','影像导航路径体现多技术来源汇入主链','dense')

# 15 conclusion
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s,NAVY)
text_box(s,'总结：这篇论文的价值不只是“画主路径”',0.75,0.8,10.8,0.5,24,WHITE,True)
text_box(s,'而是把技术演化解释推进到 claim-level technical element network。',0.78,1.45,10.6,0.32,14,CYAN,False)
card(s,'贡献','TI 衡量继承强度\nTEP 识别重要技术元素\nRNIT + Main Paths 展示演化\nOrigin analysis 解释技术融合',0.8,2.35,3.65,2.15,CYAN,13,11)
card(s,'发现','手术机器人形成四条主路径\n影像导航路径体现多技术融合\n微处理器、微创器械等后期升温',4.85,2.35,3.65,2.15,ORANGE,13,11)
card(s,'局限','0-1 词袋语义弱\n未考虑 claim 结构\ncitation time-delay\n仅限 USPTO 数据',8.9,2.35,3.65,2.15,PURPLE,13,11)
text_box(s,'最终一句话：Claims 提供技术细节，Citations 提供继承关系；二者结合，才能看清技术从哪里来、如何融合、怎样演化。',1.2,5.55,10.9,0.52,16,WHITE,True,PP_ALIGN.CENTER)
text_box(s,'15/15',12.25,7.05,0.7,0.2,8.5,RGBColor(180,204,225),True,PP_ALIGN.RIGHT)
record(15,'总结','claim-level technical element network 是本文核心价值','anchor')

# ---------- notes and source digest ----------
notes=['# Speaker Notes｜DeckForge v2\n','总时长：15 分钟。方法页不推导公式，只解释公式背后的动作。\n']
for sp in slide_specs:
    notes.append(f"## Slide {sp['page']}｜{sp['title']}\n")
    notes.append(f"- Takeaway：{sp['takeaway']}\n")
    if sp['page'] in [6,8,9]:
        notes.append('- 讲法：公式已经渲染成图片；课堂上按步骤解释，不逐符号推导。\n')
    elif sp['page']==14:
        notes.append('- 讲法：把 Path 2 讲成多技术汇入的故事，这是全 deck 的案例高潮。\n')
    else:
        notes.append('- 讲法：先读标题结论，再补页面中的 2–3 个支撑模块。\n')
(OUT/'speaker-notes.md').write_text('\n'.join(notes), encoding='utf-8')
(OUT/'slide-spec.json').write_text(json.dumps(slide_specs, ensure_ascii=False, indent=2), encoding='utf-8')
source_digest = '''# Source Digest

## Paper
Technological origination and evolution analysis by combining patent claims and citations: A case of surgical robot domain. Advanced Engineering Informatics 58 (2023) 102145.

## Core Problem
Existing technology trend analysis often shows development paths but ignores technological origins and fusion mechanisms.

## Core Method
Patent claims are treated as technical elements; patent citations are treated as technological inheritance links. The paper proposes TI, TEP, RNIT/main paths, and origin/development analysis.

## Case Facts
- Domain: surgical robot
- Data source: Patsnap / USPTO granted patents
- Patents: 3313
- Internal citations: 13097
- Average degree: 7.906
- Median TI: 0.122
- Average TI: 0.077

## Key Findings
Four main paths are identified: positioning/cutting, scanning/image-guided navigation, ophthalmic/head surgery, and bone surgery. Path 2 shows fusion of positioning, X-ray scanning, picture fusion, volumetric model, anatomical-feature positioning, virtual guidewire, and navigational guidance.
'''
(OUT/'source-digest.md').write_text(source_digest, encoding='utf-8')
image_prompts = '''# Image Prompts

No standalone AI image generation was used in v2 to avoid unnecessary image attachments. Visual assets are controlled vector diagrams and formula PNG panels. If AI images are needed later, use the Deck Style Anchor in spec_lock.json and only for cover/section background/conceptual visuals, not exact formulas or text.
'''
(OUT/'image-prompts.md').write_text(image_prompts, encoding='utf-8')

pptx = OUT/'Technological_Origination_Evolution_Surgical_Robot_DeckForge_v2.pptx'
prs.save(pptx)

# QA inspect pptx xml for common replacement chars from our generation
bad=[]
with zipfile.ZipFile(pptx,'r') as z:
    for n in z.namelist():
        if n.startswith('ppt/slides/slide') and n.endswith('.xml'):
            data=z.read(n).decode('utf-8', errors='replace')
            if '�' in data or '����' in data:
                bad.append(n)
qa = {
    'pptx': str(pptx.resolve()),
    'slides': len(slide_specs),
    'formula_policy': 'formula images used: formula_ti.png, formula_tep.png, formula_path.png',
    'density_policy': 'all content pages have >=2 modules and >=1 visual anchor by construction',
    'xml_replacement_char_scan': 'PASS' if not bad else f'WARN: {bad}',
    'standalone_image_generation': 'not used in v2 rebuild',
    'editable_elements': 'text, cards, nodes, arrows are native PPT shapes; formulas are PNG panels to prevent乱码'
}
(OUT/'qa-report.json').write_text(json.dumps(qa, ensure_ascii=False, indent=2), encoding='utf-8')
print(pptx.resolve())
print(json.dumps(qa, ensure_ascii=False, indent=2))
