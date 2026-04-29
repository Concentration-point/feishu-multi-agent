from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path
import json, zipfile, shutil

OUT = Path('deckforge_output/surgical_robot_patent_evolution_ai_generated')
OUT.mkdir(parents=True, exist_ok=True)
IMG_DIR = OUT/'ai_assets'
IMG_DIR.mkdir(exist_ok=True)

assets = {
    'cover': r'C:\Users\25723\.openclaw\media\tool-image-generation\deckforge_ai_slide_cover---4525a2d5-2097-4fe5-88a2-832c885b8c06.png',
    'method': r'C:\Users\25723\.openclaw\media\tool-image-generation\deckforge_ai_slide_method_framework---66f67e3c-21d5-4ee1-89b4-e433571225fe.png',
    'formula': r'C:\Users\25723\.openclaw\media\tool-image-generation\deckforge_ai_slide_formula_visual---64a8003d-8986-41fc-9395-10c122288172.png',
    'case': r'C:\Users\25723\.openclaw\media\tool-image-generation\deckforge_ai_slide_case_dashboard---71bd7b4f-ef64-42a0-a0b1-578592e313b3.png',
    'paths': r'C:\Users\25723\.openclaw\media\tool-image-generation\deckforge_ai_slide_paths_fusion---a1092507-ad16-41cf-abf2-70bb320d970c.png',
}
for k,p in assets.items():
    dst=IMG_DIR/f'{k}.png'
    shutil.copyfile(p,dst)
    assets[k]=str(dst)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY=RGBColor(8,18,34)
WHITE=RGBColor(255,255,255)
CYAN=RGBColor(28,190,230)
ORANGE=RGBColor(242,145,48)
GREEN=RGBColor(48,174,120)
PURPLE=RGBColor(139,92,246)
TEXT=RGBColor(25,35,50)
MUTED=RGBColor(90,105,125)
CARD=RGBColor(255,255,255)
GRID=RGBColor(218,230,242)
FONT='Microsoft YaHei'

slides=[]

def set_run(run,size=14,color=TEXT,bold=False):
    run.font.name=FONT; run.font.size=Pt(size); run.font.color.rgb=color; run.font.bold=bold

def tb(slide,text,x,y,w,h,size=14,color=TEXT,bold=False,align=None):
    box=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=box.text_frame; tf.clear(); tf.word_wrap=True; tf.auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i,line in enumerate(text.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        if align: p.alignment=align
        p.space_after=Pt(2)
        r=p.add_run(); r.text=line; set_run(r,size,color,bold)
    return box

def bg(slide,key):
    # Generated images are 3:2. Fill 16:9 slide by width, crop vertically via negative y.
    slide.shapes.add_picture(assets[key], Inches(0), Inches(-0.61), width=Inches(13.333), height=Inches(8.889))

def block(slide,x,y,w,h,fill=CARD,line=GRID):
    s=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb=fill
    s.line.color.rgb=line; s.line.width=Pt(0.8)
    return s

def title(slide,t,sub,page,dark=False):
    color=WHITE if dark else NAVY
    subc=CYAN if dark else MUTED
    tb(slide,t,0.55,0.28,8.8,0.45,21,color,True)
    if sub: tb(slide,sub,0.58,0.78,9.2,0.25,9.5,subc)
    tb(slide,f'{page:02d}/15',12.15,7.08,0.7,0.2,8.5,WHITE if dark else MUTED,True,PP_ALIGN.RIGHT)

def card(slide,t,b,x,y,w,h,accent=CYAN):
    block(slide,x,y,w,h,CARD,GRID)
    bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(0.07),Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb=accent; bar.line.fill.background()
    tb(slide,t,x+0.18,y+0.13,w-0.3,0.25,12.5,NAVY,True)
    tb(slide,b,x+0.18,y+0.48,w-0.3,h-0.55,10.2,TEXT)

def overlay_strip(slide,y=6.55):
    s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0),Inches(y),Inches(13.333),Inches(0.78))
    s.fill.solid(); s.fill.fore_color.rgb=NAVY
    s.line.fill.background()
    return s

def make_slide(key,page,t,sub,cards,foot=None,dark_title=False):
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,key); title(s,t,sub,page,dark_title)
    for c in cards:
        card(s,*c)
    if foot:
        overlay_strip(s,6.43)
        tb(s,foot,0.75,6.62,11.8,0.28,13,WHITE,True,PP_ALIGN.CENTER)
    slides.append({'page':page,'title':t,'asset':key,'note':sub})
    return s

# 1
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,'cover')
block(s,0.55,0.75,6.7,2.25,NAVY,NAVY)
tb(s,'基于 Patent Claims 与 Citations 的\n技术起源与演化分析',0.85,1.05,6.1,0.95,27,WHITE,True)
tb(s,'以手术机器人领域为例｜15 分钟课堂汇报',0.9,2.25,5.8,0.28,12,CYAN)
card(s,'核心问题','技术不仅如何发展，更重要的是：从哪里来、怎样融合。',0.8,5.75,4.4,0.82,CYAN)
tb(s,'01/15',12.15,7.08,0.7,0.2,8.5,WHITE,True,PP_ALIGN.RIGHT)
slides.append({'page':1,'title':'Cover','asset':'cover','note':'AI-generated full-slide cover'})

make_slide('cover',2,'汇报主线','Deck Brief：15 分钟只讲一条线',[
    ('Why','现有趋势分析会画路径，但较少解释技术起源。',0.75,1.25,3.55,1.0,CYAN),
    ('How','Claims = technical elements；Citations = inheritance links。',4.65,1.25,3.65,1.0,ORANGE),
    ('Case','手术机器人领域：3313 件 USPTO 授权专利。',8.65,1.25,3.55,1.0,PURPLE),
    ('So What','主路径不只是路线图，而是技术融合故事。',4.0,5.25,5.3,0.85,GREEN)
],'一句话：把 patent-level evolution 推进到 claim-level origin analysis。')

make_slide('method',3,'研究缺口','趋势分析很多，技术起源分析不足',[
    ('Common trend','统计、时间序列、机器学习：看宏观趋势，弱在技术细节。',0.65,1.15,3.35,1.25,CYAN),
    ('Patent text','关键词、LDA、SAO：能看热点，但继承关系不足。',4.25,1.15,3.35,1.25,ORANGE),
    ('Citation path','引用网络、main path：知道谁影响谁，但不知道影响了什么。',7.85,1.15,4.0,1.25,PURPLE),
    ('本文补位','把 claims 放进 citation network，追踪技术元素来源。',1.0,5.25,4.6,0.82,GREEN),
    ('核心贡献','用 TI、TEP、RNIT 和 origin analysis 形成完整框架。',6.0,5.25,5.55,0.82,CYAN)
])

make_slide('method',4,'总体框架','Source → Claims → TI → TEP → Main Paths → Origins',[
    ('1 数据收集','Patsnap 检索 surgical robot domain，采集 title、abstract、claims、citations、issue date。',0.65,1.05,3.7,1.25,CYAN),
    ('2 技术元素','每个 claim 被视为一个 technical element。',4.65,1.05,3.5,1.0,ORANGE),
    ('3 继承计算','在 citation links 上计算 claim similarity 和 TI。',8.45,1.05,3.6,1.0,PURPLE),
    ('4 重要性','TEP 衡量 claim 在时间区间内被后续技术继承的程度。',0.9,5.25,5.1,0.82,GREEN),
    ('5 演化解释','RNIT + main paths + claim origins 解释融合过程。',6.35,5.25,5.55,0.82,CYAN)
])

make_slide('method',5,'Claims 如何进入模型？','Claims 被向量化为微观技术元素',[
    ('Claims = protected invention points','比 abstract 更细，比 description 更聚焦。',0.65,1.15,4.15,1.0,CYAN),
    ('Preprocess','去除低信息词，构造技术词 Dict。',5.05,1.15,3.1,1.0,ORANGE),
    ('0–1 matrix','claim 中出现某技术词记为 1，否则为 0。',8.45,1.15,3.5,1.0,PURPLE),
    ('课堂讲法','不用陷入 NLP 细节，只讲：作者把权利要求转成可比较的技术词向量。',1.1,5.35,10.8,0.75,GREEN)
])

make_slide('formula',6,'TI：技术继承度','引用关系 j → i 上，claim-to-claim 最大相似度的平均值',[
    ('动作 1','对引用专利 i 的每个 claim，找被引专利 j 中最相似 claim。',0.7,1.05,3.65,1.12,CYAN),
    ('动作 2','每个 claim 得到一个最大相似值 MCSV。',4.7,1.05,3.45,1.0,ORANGE),
    ('动作 3','把所有 MCSV 平均，得到 TI(i,j)。',8.45,1.05,3.45,1.0,PURPLE),
    ('注意','论文剔除 TI > 0.7 的连续更新型引用，避免 novelty 弱的关系干扰。',1.1,5.45,10.7,0.78,GREEN)
])

make_slide('formula',7,'Main Technical Origin','每个 claim 在所有被引专利中找主来源',[
    ('问题','一个专利引用很多旧专利，某个 claim 的技术来源可能只对应其中一个 claim。',0.65,1.1,4.1,1.15,CYAN),
    ('MSCI','遍历所有 cited patents，取最大 MCSV 对应的 claim。',5.05,1.1,3.25,1.15,ORANGE),
    ('价值','把“引用关系”变成“技术元素来源关系”。',8.6,1.1,3.4,1.0,PURPLE),
    ('案例','US5887121 的机器人关节控制 claim 可追溯到 US5297057 / US5377310 等早期技术。',1.0,5.35,10.9,0.8,GREEN)
])

make_slide('formula',8,'TEP：技术元素持续性','一个 claim 被后续技术持续继承，TEP 就越高',[
    ('TICV','直接引用中，后续 claim 对前序 claim 贡献技术重要性。',0.7,1.05,3.55,1.05,CYAN),
    ('Indirect path','间接引用贡献沿路径传播，体现技术影响力延续。',4.55,1.05,3.55,1.05,ORANGE),
    ('TEP','在时间区间 α 内累加所有直接/间接贡献。',8.4,1.05,3.55,1.05,PURPLE),
    ('解释','TEP 高，不代表词出现多，而代表这个技术元素被后续专利持续继承。',1.05,5.45,10.9,0.78,GREEN)
])

make_slide('formula',9,'Main Paths 构建','先筛重要专利，再构建 RNIT',[
    ('Patent TEP','一个专利的 TEP = 其所有 claims 的 TEP 之和。',0.7,1.1,3.65,1.0,CYAN),
    ('Top Q','每个 5 年区间选 TEP 前 Q 的重要专利；案例中 Q = 15。',4.7,1.1,3.65,1.0,ORANGE),
    ('RNIT','搜索重要专利之间的 citation paths，形成重要技术关系网络。',8.7,1.1,3.65,1.0,PURPLE),
    ('筛选参数','路径长度至少 a = 4；路径中重要专利数量至少 b = 3。',1.2,5.35,10.75,0.78,GREEN)
])

make_slide('case',10,'案例数据','Surgical Robot Domain：3313 件专利，13097 条内部引用',[
    ('3313','USPTO granted patents',0.65,1.1,2.1,0.9,CYAN),
    ('13097','internal citations',3.0,1.1,2.1,0.9,ORANGE),
    ('7.906','average degree',5.35,1.1,2.1,0.9,PURPLE),
    ('0.122','median TI',7.7,1.1,2.1,0.9,GREEN),
    ('0.077','average TI',10.05,1.1,2.1,0.9,CYAN),
    ('TI 分布','Long-tail：多数引用不是强技术继承，因此不能把 citation 一视同仁。',1.05,5.45,10.9,0.78,ORANGE)
])

make_slide('case',11,'技术热点变化','TEP 展示 2011–2015 与 2016–2020 的重要技术元素变化',[
    ('2011–2015','robotic arm motion control；X-ray imaging；endoscope control；bone positioning。',0.75,1.15,5.3,1.25,CYAN),
    ('2016–2020','surgical operation systems；microprocessor；endoscope insertion；minimally invasive robots；stapling device。',6.5,1.15,5.65,1.25,ORANGE),
    ('变化 1','US5236432 相关技术排名从第 5 升到第 2。',0.9,5.25,3.6,0.82,GREEN),
    ('变化 2','microprocessor 和 stapling technique 可能是后续机会。',4.85,5.25,3.85,0.82,PURPLE),
    ('变化 3','技术中心从机械/影像控制扩展到系统集成和微创器械。',9.05,5.25,3.3,0.82,CYAN)
])

make_slide('paths',12,'四条主路径','手术机器人技术演化的四个子领域故事',[
    ('Path 1｜定位与切割','laser surgical knife → laser surgery → robot-aided surgery → robotic tools/data architecture',0.7,1.05,5.65,1.05,CYAN),
    ('Path 2｜扫描/影像导航','stereotactic apparatus → CT-assisted surgery → image-directed robot → navigational guidance',6.75,1.05,5.65,1.05,ORANGE),
    ('Path 3｜眼科/头部','fluid control → intraocular suction → ophthalmic instruments → head probe positioning',0.7,5.2,5.65,1.05,GREEN),
    ('Path 4｜骨科手术','bone alteration → image-directed robot → skeletal surgery → fluoroscopic navigation',6.75,5.2,5.65,1.05,PURPLE)
])

make_slide('paths',13,'重点路径：Path 2 技术融合','扫描/影像导航路径最能体现 origin analysis 的解释力',[
    ('起点','US4341220：stereotactic surgery apparatus，融合 positioning 与 X-ray scanning。',0.75,1.05,4.0,1.05,CYAN),
    ('中段','US4791934 / US5086401：CT、picture fusion、digital imaging 与 bone positioning 汇入。',5.05,1.05,4.0,1.05,ORANGE),
    ('后段','US5682886 / US6470207：volumetric model、3D prosthesis imaging、virtual guidewire、navigation。',9.35,1.05,3.25,1.25,PURPLE),
    ('一句话','Path 2 不是单线进步，而是定位、影像、三维建模和导航的连续融合。',1.15,5.45,10.85,0.78,GREEN)
])

make_slide('paths',14,'论文评价','贡献清楚，但 NLP 表征偏弱',[
    ('Strength 1','从 claims 层面分析，比 title / abstract 更接近创新点。',0.75,1.1,3.8,1.0,CYAN),
    ('Strength 2','TI / TEP / main paths 形成可计算框架。',4.85,1.1,3.45,1.0,ORANGE),
    ('Strength 3','claim-level origin analysis 能讲清技术融合故事。',8.6,1.1,3.7,1.0,GREEN),
    ('Limitation','0-1 词袋 + cosine similarity 语义弱；未考虑 claim 结构；citation 有时间滞后；数据仅限 USPTO。',1.05,5.25,11.1,0.9,PURPLE)
])

s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s,'cover')
block(s,0.85,0.9,11.65,1.15,NAVY,NAVY)
tb(s,'最终结论',1.15,1.12,2.2,0.35,23,WHITE,True)
tb(s,'Claims 提供技术细节，Citations 提供继承关系；二者结合，才能看清技术从哪里来、如何融合、怎样演化。',3.15,1.16,8.85,0.3,13,CYAN,True)
card(s,'方法','TI 衡量继承强度\nTEP 衡量技术元素重要性\nRNIT / Main Paths 展示演化路径',1.0,3.0,3.45,1.8,CYAN)
card(s,'案例','手术机器人识别四条主路径\nPath 2 体现影像导航类技术融合',4.95,3.0,3.45,1.8,ORANGE)
card(s,'启示','专利分析应从“引用网络”升级到“技术元素网络”。',8.9,3.0,3.45,1.8,PURPLE)
tb(s,'15/15',12.15,7.08,0.7,0.2,8.5,WHITE,True,PP_ALIGN.RIGHT)
slides.append({'page':15,'title':'Conclusion','asset':'cover','note':'AI visual conclusion'})

# notes and metadata
(OUT/'slide-spec.json').write_text(json.dumps(slides, ensure_ascii=False, indent=2), encoding='utf-8')
(OUT/'style-bible.json').write_text(json.dumps({
    'approach':'AI-generated full-slide visual bases + editable text overlays for factual accuracy',
    'assets':assets,
    'warning':'Generated images intentionally contain no readable text to avoid hallucinated/garbled text; all real text is native PPT overlay.'
}, ensure_ascii=False, indent=2), encoding='utf-8')
notes=['# Speaker Notes｜AI-generated visual deck\n','这版使用生图模型生成整页视觉底图；准确中文、数字和公式解释用 PPT 原生文本覆盖，避免模型乱写字。\n']
for s in slides:
    notes.append(f"## {s['page']:02d}. {s['title']}\n- Visual asset: {s['asset']}\n- 讲法：先读标题结论，再用页面卡片支撑。\n")
(OUT/'speaker-notes.md').write_text('\n'.join(notes), encoding='utf-8')
qa={
    'slides':15,
    'image_generation':'5 full-slide AI visual assets generated with openai/gpt-image-2 and reused consistently',
    'text_policy':'real Chinese text and facts are editable overlays, not generated inside images',
    'reason':'image models are unreliable for exact Chinese/formulas; this preserves visual quality plus factual correctness'
}
(OUT/'qa-report.json').write_text(json.dumps(qa, ensure_ascii=False, indent=2), encoding='utf-8')

pptx=OUT/'Technological_Origination_Evolution_Surgical_Robot_AI_Image_Deck_v3.pptx'
prs.save(pptx)
with zipfile.ZipFile(pptx,'r') as z:
    bad=[]
    for n in z.namelist():
        if n.startswith('ppt/slides/slide') and n.endswith('.xml'):
            d=z.read(n).decode('utf-8',errors='replace')
            if '�' in d: bad.append(n)
qa['xml_replacement_char_scan']='PASS' if not bad else str(bad)
(OUT/'qa-report.json').write_text(json.dumps(qa, ensure_ascii=False, indent=2), encoding='utf-8')
print(pptx.resolve())
print(json.dumps(qa, ensure_ascii=False, indent=2))
