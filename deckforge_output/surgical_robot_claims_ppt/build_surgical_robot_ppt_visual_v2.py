from pathlib import Path
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE = Path(__file__).resolve().parent
FIG = BASE / 'assets' / 'paper_figures'
OUT = BASE / 'output' / 'final' / 'surgical_robot_claims_calm_academic_visual_v2.pptx'
W, H = Inches(13.333333), Inches(7.5)

COL = {
    'bg': 'F7F8FA', 'panel': 'FFFFFF', 'primary': '1F4E79', 'accent': 'D89A2B',
    'text': '222222', 'muted': '6B7280', 'dark': '0F172A', 'soft': 'E8EEF5',
    'line': 'D9E2EC', 'pale': 'EEF4FA'
}
FONT = 'Microsoft YaHei'


def rgb(h):
    h = h.strip('#')
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:], 16))


def bg(slide, color='bg'):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = rgb(COL[color])


def txt(slide, x, y, w, h, s, size=16, bold=False, color='text', align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(0.03); tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.text = s
    if align is not None: p.alignment = align
    r = p.runs[0]
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(COL[color])
    return box


def bullet(slide, x, y, w, h, items, size=14, color='text'):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(item); p.level = 0
        p.font.name = FONT; p.font.size = Pt(size); p.font.color.rgb = rgb(COL[color])
        p.space_after = Pt(5)
    return box


def rect(slide, x, y, w, h, fill='panel', line=None, radius=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(COL[fill])
    shape.line.color.rgb = rgb(COL[line or fill])
    return shape


def line(slide, x1, y1, x2, y2, color='accent', width=1.5):
    shp = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shp.line.color.rgb = rgb(COL[color]); shp.line.width = Pt(width)
    return shp


def add_img(slide, path, x, y, w, h, border=True):
    path = Path(path)
    if not path.exists():
        rect(slide, x, y, w, h, 'soft', 'line')
        txt(slide, x+0.1, y+0.1, w-0.2, h-0.2, f'Missing image:\n{path.name}', 10, False, 'muted')
        return
    if border:
        rect(slide, x-0.03, y-0.03, w+0.06, h+0.06, 'panel', 'line')
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def title(slide, s, sub=None, idx=None):
    txt(slide, 0.65, 0.38, 10.6, 0.38, s, 22, True, 'primary')
    if sub: txt(slide, 0.67, 0.82, 10.4, 0.26, sub, 9, False, 'muted')
    if idx: txt(slide, 12.15, 7.02, 0.55, 0.2, f'{idx:02d}', 9, False, 'muted', PP_ALIGN.RIGHT)


def pill(slide, x, y, label, fill='pale', color='primary'):
    rect(slide, x, y, 1.55, 0.32, fill, fill)
    txt(slide, x+0.08, y+0.07, 1.38, 0.13, label, 8, True, color, PP_ALIGN.CENTER)


def make_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    rect(s, 0.55, 0.5, 12.25, 6.45, 'panel', 'line')
    txt(s, 1.05, 1.05, 3.2, 0.28, '课堂汇报 · Calm Academic', 11, True, 'accent')
    txt(s, 1.05, 1.8, 6.7, 1.05, '专利权利要求 + 引文\n技术起源与演化分析', 30, True, 'primary')
    txt(s, 1.08, 3.15, 6.0, 0.48, 'Technological origination and evolution analysis by combining patent claims and citations', 13, False, 'muted')
    bullet(s, 1.1, 4.25, 5.8, 0.9, ['Qiu & Wang · Advanced Engineering Informatics · 2023', 'Case: surgical robot domain · 3313 USPTO patents'], 12)
    # visual motif: citation/claim network
    for i, (x, y, lab) in enumerate([(8.2,1.6,'Claim'), (10.0,2.3,'Citation'), (8.7,3.6,'TI'), (10.7,4.5,'Path')]):
        rect(s, x, y, 1.45, 0.72, 'pale', 'line')
        txt(s, x+0.12, y+0.23, 1.2, 0.16, lab, 11, True, 'primary', PP_ALIGN.CENTER)
    line(s, 9.55,1.95,10.0,2.55); line(s,10.7,2.95,9.4,3.75); line(s,9.6,4.0,10.7,4.65)


def make_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s, '研究问题：趋势分析还不够，还要解释技术从哪来', '核心缺口：现有方法重发展轨迹，轻技术 origin 与融合机制', 2)
    rect(s, 0.8, 1.35, 5.4, 4.95, 'panel', 'line')
    txt(s, 1.08, 1.68, 4.8, 0.25, '传统 patent trend analysis', 17, True, 'primary')
    bullet(s, 1.12, 2.2, 4.65, 1.6, ['统计 / 时间序列', '文本挖掘：keyword、SAO、LDA', 'citation network / main path'], 13)
    rect(s, 1.1, 4.25, 4.65, 1.25, 'soft', 'line')
    txt(s, 1.35, 4.55, 4.1, 0.28, '问题：知道“怎么发展”，但很难知道“技术细节从哪继承”。', 14, True, 'dark')
    rect(s, 6.9, 1.35, 5.65, 4.95, 'panel', 'line')
    txt(s, 7.18, 1.68, 4.8, 0.25, '本文的切入点', 17, True, 'primary')
    add_img(s, FIG/'fig1_workflow_crop.jpg', 7.25, 2.12, 4.9, 2.55)
    txt(s, 7.25, 5.05, 4.9, 0.42, '把 claim 当作 technical element，再用 citation 提供继承方向。', 13, True, 'accent', PP_ALIGN.CENTER)


def make_framework(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s, '方法总览：claim-level 技术元素 + citation 继承方向', '这页用信息图承载框架，减少文本框拥挤', 3)
    add_img(s, BASE / 'assets' / 'framework_pipeline.png', 0.85, 1.35, 11.65, 4.85)
    txt(s, 1.0, 6.45, 11.2, 0.25, '一句话：claims 给“技术细节”，citations 给“继承链路”，main path 给“演化故事”。', 13, True, 'accent', PP_ALIGN.CENTER)

def make_ti(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s, '关键指标 TI：把 citation 从“有无关系”变成“继承强度”', '公式改成安全示意，避免 PPT 数学符号乱码', 4)
    # Native diagram instead of formulas
    rect(s, 0.9, 1.5, 3.0, 1.2, 'panel', 'line'); txt(s,1.15,1.82,2.5,0.22,'Citing patent i',15,True,'primary',PP_ALIGN.CENTER); txt(s,1.15,2.18,2.5,0.18,'claim_i_1 ... claim_i_m',9,False,'muted',PP_ALIGN.CENTER)
    rect(s, 9.45, 1.5, 3.0, 1.2, 'panel', 'line'); txt(s,9.7,1.82,2.5,0.22,'Cited patent j',15,True,'primary',PP_ALIGN.CENTER); txt(s,9.7,2.18,2.5,0.18,'claim_j_1 ... claim_j_k',9,False,'muted',PP_ALIGN.CENTER)
    line(s,3.95,2.1,9.35,2.1); txt(s,5.0,1.72,3.2,0.22,'max cosine similarity',13,True,'accent',PP_ALIGN.CENTER)
    rect(s, 2.2, 3.55, 8.9, 1.55, 'soft', 'line')
    txt(s, 2.55, 3.85, 8.2, 0.25, 'TI(i,j) = average of max claim similarities', 17, True, 'dark', PP_ALIGN.CENTER)
    txt(s, 2.7, 4.38, 7.9, 0.28, 'Interpretation: higher TI means stronger claim-level inheritance in a citation.', 11, False, 'muted', PP_ALIGN.CENTER)
    bullet(s, 1.3, 5.75, 10.6, 0.55, ['不在 PPT 文本框里硬放复杂下标公式；课堂讲解用“平均最大相似度”即可。'], 13)


def make_origin_example(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s, '例子：为每个 claim 找 main technical origin', '文献图替代纯文字解释：US5887121 的 claim 来源分析', 5)
    add_img(s, FIG/'fig17_origin_crop.jpg', 0.75, 1.28, 7.1, 5.2)
    rect(s, 8.15, 1.45, 4.35, 4.85, 'panel', 'line')
    txt(s, 8.45, 1.8, 3.75, 0.25, '怎么读这页图？', 16, True, 'primary')
    bullet(s, 8.5, 2.35, 3.5, 2.5, ['右侧：目标专利 claims', '左侧：被引用专利中的来源 claims', '箭头：citation + overlap terms', '用途：解释技术元素如何继承'], 12)
    txt(s, 8.5, 5.45, 3.5, 0.35, 'Takeaway：origin analysis 能把抽象 citation 拆成具体 claim 来源。', 12, True, 'accent')


def make_tep(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s, 'TEP：按时间窗找重要技术元素', '用文献 Table 4 展示 2011–2015 与 2016–2020 的变化', 6)
    add_img(s, FIG/'table4_important_elements_crop.jpg', 0.75, 1.25, 7.25, 5.15)
    rect(s, 8.3, 1.45, 4.1, 4.75, 'panel', 'line')
    txt(s, 8.6, 1.78, 3.5, 0.25, '变化趋势', 16, True, 'primary')
    bullet(s, 8.65, 2.28, 3.25, 2.8, ['2011–2015：机械臂运动控制、X-ray / CT、内窥镜控制', '2016–2020：系统、微处理器、微创与内窥镜更突出', 'TEP 快速增大：说明领域近年发展加速'], 12)
    txt(s, 8.65, 5.5, 3.25, 0.3, '别堆全表，讲排名迁移。', 12, True, 'accent')


def make_case_data(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s, '案例数据与网络：3313 件 surgical robot 专利', '数据卡 + 文献网络图，降低全文字感', 7)
    metrics=[('3313','patents'),('13,097','internal citations'),('7.906','avg degree'),('0.122','median TI'),('0.077','avg TI')]
    x=0.72
    for val,lab in metrics:
        rect(s,x,1.35,2.2,0.95,'panel','line')
        txt(s,x+0.1,1.55,2.0,0.25,val,18,True,'primary',PP_ALIGN.CENTER)
        txt(s,x+0.1,1.9,2.0,0.18,lab,8,False,'muted',PP_ALIGN.CENTER)
        x+=2.48
    add_img(s, FIG/'fig21_rnit_crop.jpg', 0.85, 2.65, 6.0, 3.65)
    add_img(s, FIG/'data_stats_crop.jpg', 7.25, 2.65, 5.1, 3.65)


def make_paths(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s, '四条主路径：把复杂网络压成可讲的演化路线', '用文献 Table 5 + 右侧提炼，而不是满页文字', 8)
    add_img(s, FIG/'table5_main_paths_crop.jpg', 0.75, 1.25, 6.65, 5.35)
    paths=[('1','定位与切割'),('2','X-ray / scanning'),('3','眼科 / 头部手术'),('4','骨科定位操作')]
    y=1.55
    for no,lab in paths:
        rect(s,7.75,y,4.45,0.72,'panel','line')
        txt(s,8.0,y+0.18,0.45,0.15,no,13,True,'accent',PP_ALIGN.CENTER)
        txt(s,8.65,y+0.17,3.1,0.18,lab,14,True,'primary')
        y+=0.95
    rect(s,7.75,5.55,4.45,0.65,'soft','line')
    txt(s,8.05,5.78,3.8,0.15,'Takeaway：main path 是演化叙事，不是单纯节点列表。',10,True,'dark',PP_ALIGN.CENTER)


def make_fusion(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s, '技术融合：定位 + 影像 + 导航逐步组合', '文献 Fig.23 是这篇论文最值得展示的图', 9)
    add_img(s, FIG/'fig23_fusion_process_crop.jpg', 0.65, 1.18, 8.1, 5.75)
    rect(s,9.05,1.4,3.45,4.95,'panel','line')
    txt(s,9.3,1.72,3.0,0.25,'Path 2 的演化逻辑',15,True,'primary')
    bullet(s,9.35,2.2,2.75,2.8,['Stereotactic apparatus', 'CT-assisted robotic system', 'Image-directed surgery', '3D / volumetric model', 'Real-time navigation'],11)
    txt(s,9.35,5.45,2.75,0.35,'这是“技术融合”最直观的证据页。',11,True,'accent')


def make_conclusion(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s, '结论：claim-level 视角让技术演化更可解释', '最后用三列收束，不再塞长段落', 10)
    cols=[('贡献',['TI：claim-based inheritance','TEP：technical element persistence','Main path：技术细节路径','Origin analysis：融合来源']),('发现',['未来热点：数字终端 / 内窥镜 / 微创机器人','四条路线：切割、扫描、头部、骨科','融合主线：定位 + 影像 + 导航']),('局限',['未考虑 claim 结构','citation 有时间滞后','cosine similarity 偏浅','可用更强 NLP 改进'])]
    x=0.85
    for head,items in cols:
        rect(s,x,1.55,3.65,4.6,'panel','line')
        txt(s,x+0.25,1.9,3.1,0.25,head,17,True,'primary')
        bullet(s,x+0.3,2.45,2.9,2.7,items,12)
        x+=4.0
    txt(s,1.0,6.55,11.3,0.25,'课堂一句话：这篇论文把“专利引用网络”推进到“claim 级技术继承解释”。',13,True,'accent',PP_ALIGN.CENTER)


def build():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs=Presentation(); prs.slide_width=W; prs.slide_height=H
    for f in [make_cover, make_problem, make_framework, make_ti, make_origin_example, make_tep, make_case_data, make_paths, make_fusion, make_conclusion]:
        f(prs)
    prs.save(OUT)
    print(f'Wrote {OUT}')

if __name__=='__main__':
    build()
